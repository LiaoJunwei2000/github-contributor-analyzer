"""华为浅色版 PPT 生成器（16:9，参考华为官方浅色模板）。

配色参考：
  - 华为红   #C7000B  —— 主色，页眉/强调
  - 深灰     #333333  —— 主体文字
  - 中灰     #6B6B6B  —— 辅助文字
  - 浅灰     #F5F5F5  —— 卡片/区块底色
  - 边框灰   #E0E0E0  —— 细线/分割
  - 白色     #FFFFFF  —— 背景
"""

import io
import os
import re
import math
import datetime
import urllib.request
import traceback

import matplotlib
import matplotlib.font_manager as _fm

# 尝试注册系统中文字体（macOS / Linux / Windows 常见路径）
_CJK_CANDIDATES = [
    # macOS
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/Supplemental/Songti.ttc",
    # Linux
    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    # Windows
    "C:/Windows/Fonts/msyh.ttc",
    "C:/Windows/Fonts/simsun.ttc",
]

_CJK_FONT_NAME = "Arial"   # 默认回退
for _p in _CJK_CANDIDATES:
    if os.path.exists(_p):
        try:
            _fm.fontManager.addfont(_p)
            _prop = _fm.FontProperties(fname=_p)
            _CJK_FONT_NAME = _prop.get_name()
            break
        except Exception:
            pass

matplotlib.rcParams.update({
    "font.sans-serif": [_CJK_FONT_NAME, "Arial", "DejaVu Sans"],
    "axes.unicode_minus": False,
})
import matplotlib.pyplot as plt
import numpy as np
import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from db import list_repos, get_contributors, list_tags, get_repos_by_tags, get_all_repo_tags

# ── 华为浅色版配色 ────────────────────────────────────────────
C_RED    = RGBColor(0xC7, 0x00, 0x0B)   # 华为红（主色）
C_DRED   = RGBColor(0xA3, 0x00, 0x09)   # 深红（强调/hover）
C_DARKRED= RGBColor(0x6B, 0x00, 0x04)   # 暗红（替代深蓝，封面/对比色）
C_DGRAY  = RGBColor(0x33, 0x33, 0x33)   # 主文字
C_MGRAY  = RGBColor(0x6B, 0x6B, 0x6B)   # 辅助文字
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)   # 卡片/区块底
C_BORDER = RGBColor(0xE0, 0xE0, 0xE0)   # 细线/边框
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x00, 0x7A, 0x3D)   # 求职绿

# ── 尺寸 ─────────────────────────────────────────────────────
SW = Cm(33.87)          # 16:9 宽
SH = Cm(19.05)          # 16:9 高

# 内容页通用区域
HDR_H  = Cm(2.6)        # 页眉高
FTR_H  = Cm(0.9)        # 页脚高
BODY_T = HDR_H          # 内容起始 Y
BODY_B = SH - FTR_H     # 内容结束 Y
BODY_H = BODY_B - BODY_T


# ════════════════════════════════════════════════════════════
# 低级绘图辅助
# ════════════════════════════════════════════════════════════

def _rect(sl, l, t, w, h, fill=None, line_rgb=None, lw_pt=0.5):
    """添加矩形色块。"""
    s = sl.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(lw_pt)
    else:
        s.line.fill.background()
    return s


def _txt(sl, text, l, t, w, h,
         sz=11, bold=False, italic=False,
         color=C_DGRAY, align=PP_ALIGN.LEFT, wrap=True, url=None):
    """添加文字框。url 不为空时将文字设为可点击超链接。"""
    tb = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text) if text not in (None, "None") else ""
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Arial"
    r.font.color.rgb = color
    if url:
        r.hyperlink.address = url
        r.font.underline = True
    return tb


def _fig_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _fetch_avatar(url: str) -> io.BytesIO | None:
    if not url:
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "PPT-Generator/1.0"})
        with urllib.request.urlopen(req, timeout=5) as r:
            buf = io.BytesIO(r.read())
            buf.seek(0)
            return buf
    except Exception:
        return None


# ════════════════════════════════════════════════════════════
# 华为浅色版：通用页眉 / 页脚
# ════════════════════════════════════════════════════════════

def _header(sl, title: str, sub: str = ""):
    """
    浅色版页眉：
      - 白色背景
      - 左侧 0.35cm 华为红竖条
      - 底部 0.08cm 浅灰分割线
      - 标题（深灰粗体）+ 副标题（中灰）
    """
    # 白色页眉底
    _rect(sl, 0, 0, SW, HDR_H, fill=C_WHITE)
    # 左侧红竖条
    _rect(sl, 0, 0, Cm(0.35), HDR_H, fill=C_RED)
    # 底部分割线
    _rect(sl, 0, HDR_H - Cm(0.08), SW, Cm(0.08), fill=C_BORDER)

    # 主标题
    _txt(sl, title, Cm(0.75), Cm(0.4), SW - Cm(2), Cm(1.3),
         sz=20, bold=True, color=C_DGRAY)
    # 副标题
    if sub:
        _txt(sl, sub, Cm(0.75), Cm(1.75), SW - Cm(2), Cm(0.72),
             sz=10, color=C_MGRAY)


def _footer(sl, repo: str = "", note: str = ""):
    """
    浅色版页脚：
      - 浅灰底条
      - 顶部 0.08cm 红色细线（华为品牌线）
      - 仓库名（左）+ 页码（右）
    """
    # 页脚底条
    _rect(sl, 0, BODY_B, SW, FTR_H, fill=C_LGRAY)
    # 顶部红细线
    _rect(sl, 0, BODY_B, SW, Cm(0.08), fill=C_RED)

    if repo:
        _txt(sl, repo, Cm(0.75), BODY_B + Cm(0.15), Cm(20), FTR_H - Cm(0.2),
             sz=8, color=C_MGRAY)
    if note:
        _txt(sl, note, SW - Cm(5.5), BODY_B + Cm(0.15), Cm(5.2), FTR_H - Cm(0.2),
             sz=8, color=C_MGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# Slide 1：封面
# 华为封面风格：左侧深色面板 + 右侧白色内容区
# ════════════════════════════════════════════════════════════

def _slide_cover(prs, repo: str, stats: dict):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 左侧深色面板（约 42% 宽）──
    LP_W = Cm(14.2)
    _rect(sl, 0, 0, LP_W, SH, fill=C_DARKRED)

    # 左侧顶部红色短条（品牌线）
    _rect(sl, 0, 0, LP_W, Cm(0.5), fill=C_RED)
    # 左侧底部红色短条
    _rect(sl, 0, SH - Cm(0.5), LP_W, Cm(0.5), fill=C_RED)

    # 左侧内容
    _txt(sl, "贡献者分析报告",
         Cm(1.2), Cm(3.8), LP_W - Cm(1.8), Cm(1.3),
         sz=14, color=RGBColor(0xFF, 0xCC, 0xCC))        # 浅红辅助标题

    _txt(sl, repo,
         Cm(1.2), Cm(5.2), LP_W - Cm(1.8), Cm(3.5),
         sz=24, bold=True, color=C_WHITE, wrap=True)

    # 分割线（白色半透明感，用浅色实现）
    _rect(sl, Cm(1.2), Cm(9.0), Cm(5), Cm(0.06),
          fill=RGBColor(0x88, 0x99, 0xBB))

    # 统计摘要（左侧下方）
    for i, (val, label) in enumerate([
        (f"{stats.get('total', 0):,}",   "贡献者总数"),
        (f"{stats.get('commits', 0):,}", "总 Commits"),
    ]):
        y = Cm(9.5) + i * Cm(2.4)
        _txt(sl, val, Cm(1.2), y, LP_W - Cm(1.8), Cm(1.4),
             sz=26, bold=True, color=C_WHITE)
        _txt(sl, label, Cm(1.2), y + Cm(1.4), LP_W - Cm(1.8), Cm(0.7),
             sz=9, color=RGBColor(0xBB, 0xCC, 0xDD))

    # ── 右侧白色内容区 ──
    RP_L = LP_W + Cm(0.08)   # 中间留细缝
    RP_W = SW - RP_L
    _rect(sl, LP_W, 0, Cm(0.08), SH, fill=C_RED)   # 红色分界线
    _rect(sl, RP_L, 0, RP_W, SH, fill=C_WHITE)

    # 右侧大标题区
    _txt(sl, "GitHub",
         RP_L + Cm(2), Cm(3.5), RP_W - Cm(2.5), Cm(1.5),
         sz=36, bold=True, color=C_RED)
    _txt(sl, "Contributor Analyzer",
         RP_L + Cm(2), Cm(5.1), RP_W - Cm(2.5), Cm(1.2),
         sz=20, bold=False, color=C_DGRAY)

    # 右侧日期
    _txt(sl, datetime.date.today().strftime("%Y  /  %m  /  %d"),
         RP_L + Cm(2), Cm(6.5), RP_W - Cm(2.5), Cm(0.8),
         sz=11, color=C_MGRAY)

    # 右侧底部灰色区域（装饰）
    _rect(sl, RP_L, SH - Cm(3.5), RP_W, Cm(3.5), fill=C_LGRAY)
    _txt(sl, "Powered by GitHub REST API",
         RP_L + Cm(2), SH - Cm(2.8), RP_W - Cm(3), Cm(0.8),
         sz=9, color=C_MGRAY)
    _txt(sl, "github.com/" + repo,
         RP_L + Cm(2), SH - Cm(1.9), RP_W - Cm(3), Cm(0.8),
         sz=10, bold=True, color=C_RED)


# ════════════════════════════════════════════════════════════
# Slide 2：公司 & 地区概览
# ════════════════════════════════════════════════════════════

def _company_chart(df: pd.DataFrame) -> io.BytesIO:
    cdf = df[df["company"].notna() & (df["company"] != "")].copy()
    cdf["company"] = cdf["company"].str.strip().str.lstrip("@")
    counts = cdf.groupby("company").size().sort_values(ascending=True).tail(12)

    if counts.empty:
        fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor="white")
        ax.text(0.5, 0.5, "No Data", ha="center", va="center",
                fontsize=11, color="#999", transform=ax.transAxes)
        ax.axis("off")
        return _fig_buf(fig)

    fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor="white")
    bars = ax.barh(counts.index, counts.values,
                   color="#C7000B", edgecolor="none", height=0.6)
    ax.bar_label(bars, fmt="%d", padding=4, fontsize=8, color="#333")
    ax.set_xlabel("贡献者人数", fontsize=8, color="#6B6B6B")
    ax.tick_params(labelsize=8, colors="#6B6B6B")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#E0E0E0")
    ax.set_facecolor("white")
    ax.grid(axis="x", linestyle="--", alpha=0.3, color="#E0E0E0")
    fig.tight_layout(pad=0.8)
    return _fig_buf(fig)


def _region_chart(df: pd.DataFrame) -> io.BytesIO:
    ldf = df[df["location"].notna() & (df["location"] != "")].copy()

    def extract(loc):
        parts = str(loc).split(",")
        r = parts[-1].strip()
        if re.match(r"^\d[\d\s\-]*$", r) or r == "":
            r = parts[-2].strip() if len(parts) >= 2 else r
        return r

    if ldf.empty:
        fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor="white")
        ax.text(0.5, 0.5, "No Data", ha="center", va="center",
                fontsize=11, color="#999", transform=ax.transAxes)
        ax.axis("off")
        return _fig_buf(fig)

    ldf["region"] = ldf["location"].apply(extract)
    counts = (
        ldf[ldf["region"] != ""].groupby("region").size()
        .sort_values(ascending=True).tail(12)
    )

    fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor="white")
    bars = ax.barh(counts.index, counts.values,
                   color="#6B0004", edgecolor="none", height=0.6)
    ax.bar_label(bars, fmt="%d", padding=4, fontsize=8, color="#333")
    ax.set_xlabel("贡献者人数", fontsize=8, color="#6B6B6B")
    ax.tick_params(labelsize=8, colors="#6B6B6B")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#E0E0E0")
    ax.set_facecolor("white")
    ax.grid(axis="x", linestyle="--", alpha=0.3, color="#E0E0E0")
    fig.tight_layout(pad=0.8)
    return _fig_buf(fig)


def _stat_card(sl, l, t, w, h, val: str, label: str):
    """华为浅色版统计卡片：白底 + 左红竖条 + 数字 + 标签。"""
    _rect(sl, int(l), int(t), int(w), int(h), fill=C_WHITE,
          line_rgb=C_BORDER, lw_pt=0.6)
    # 左红竖条
    _rect(sl, int(l), int(t), int(Cm(0.2)), int(h), fill=C_RED)
    # 数字
    _txt(sl, val, int(l + Cm(0.5)), int(t + Cm(0.15)), int(w - Cm(0.7)), int(h * 0.58),
         sz=17, bold=True, color=C_RED, align=PP_ALIGN.LEFT)
    # 标签
    _txt(sl, label, int(l + Cm(0.5)), int(t + Cm(0.15) + int(h * 0.58)),
         int(w - Cm(0.7)), int(h * 0.38),
         sz=8.5, color=C_MGRAY, align=PP_ALIGN.LEFT)


def _slide_overview(prs, df: pd.DataFrame, repo: str, pg: int, total: int):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # 白色背景
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)
    _header(sl, "贡献者概览", "来源公司 & 地区分布分析")
    _footer(sl, repo, f"{pg} / {total}")

    # ── 4 个统计卡片（横排）──
    n_total   = len(df)
    n_company = int(df["company"].notna().sum())
    n_loc     = int(df["location"].notna().sum())
    n_commits = int(df["total_commits"].sum()) if "total_commits" in df.columns else 0

    card_metas = [
        (f"{n_total:,}",   "贡献者总数"),
        (f"{n_company:,}", "有公司信息"),
        (f"{n_loc:,}",     "有地区信息"),
        (f"{n_commits:,}", "总 Commits"),
    ]
    CARD_W = (SW - Cm(1.2)) / 4
    CARD_H = Cm(1.85)
    CARD_T = BODY_T + Cm(0.4)
    for i, (val, label) in enumerate(card_metas):
        _stat_card(sl, Cm(0.6) + i * CARD_W, CARD_T, CARD_W - Cm(0.2), CARD_H, val, label)

    # ── 两张图表 ──
    CHART_T = CARD_T + CARD_H + Cm(0.55)
    CHART_H = BODY_B - CHART_T - Cm(0.2)
    MID     = SW / 2

    # 左：公司（红色条形）
    _txt(sl, "来源公司分布（Top 12）",
         Cm(0.6), CHART_T - Cm(0.5), MID - Cm(0.8), Cm(0.45),
         sz=10, bold=True, color=C_DGRAY)
    # 红色小方块标记
    _rect(sl, int(Cm(0.6)), int(CHART_T - Cm(0.42)), int(Cm(0.22)), int(Cm(0.22)), fill=C_RED)

    sl.shapes.add_picture(_company_chart(df),
                          int(Cm(0.5)), int(CHART_T),
                          int(MID - Cm(1.1)), int(CHART_H))

    # 竖分割线
    _rect(sl, int(MID), int(CHART_T - Cm(0.5)), int(Cm(0.06)),
          int(CHART_H + Cm(0.5)), fill=C_BORDER)

    # 右：地区（深蓝条形）
    _txt(sl, "来源地区分布（Top 12）",
         MID + Cm(0.8), CHART_T - Cm(0.5), MID - Cm(1.2), Cm(0.45),
         sz=10, bold=True, color=C_DGRAY)
    _rect(sl, int(MID + Cm(0.8)), int(CHART_T - Cm(0.42)),
          int(Cm(0.22)), int(Cm(0.22)), fill=C_DARKRED)

    sl.shapes.add_picture(_region_chart(df),
                          int(MID + Cm(0.6)), int(CHART_T),
                          int(MID - Cm(1.1)), int(CHART_H))


# ════════════════════════════════════════════════════════════
# Slide 3…：精简名片汇总（一页或多页，每页最多 PER_PAGE 人）
# ════════════════════════════════════════════════════════════

# 名片汇总页固定布局：2 行 × 3 列，每页 6 人
_SUMMARY_COLS     = 3
_SUMMARY_ROWS     = 2
_SUMMARY_PER_PAGE = _SUMMARY_COLS * _SUMMARY_ROWS   # = 6

# 固定卡片尺寸（根据 16:9 幻灯片精确计算）
_PAD_X    = Cm(0.55)
_PAD_Y    = Cm(0.45)
_GAP_X    = Cm(0.28)
_GAP_Y    = Cm(0.28)
_CARD_W   = (SW - _PAD_X * 2 - _GAP_X * (_SUMMARY_COLS - 1)) / _SUMMARY_COLS
_CARD_H   = (BODY_H - _PAD_Y * 2 - _GAP_Y * (_SUMMARY_ROWS - 1)) / _SUMMARY_ROWS


def _name_card(sl, cx, cy, row: dict, av_cache: dict):
    """
    固定尺寸名片（_CARD_W × _CARD_H）。
    布局：左红竖条 | 头像（居中对齐）| 文字区（姓名/login/信息）
    """
    cw, ch = _CARD_W, _CARD_H

    # 卡片底色 + 边框
    _rect(sl, int(cx), int(cy), int(cw), int(ch),
          fill=C_WHITE, line_rgb=C_BORDER, lw_pt=0.5)
    # 左红竖条
    RED_BAR = Cm(0.2)
    _rect(sl, int(cx), int(cy), int(RED_BAR), int(ch), fill=C_RED)

    login = str(row.get("login") or "?")
    name  = str(row.get("name") or login)
    rank  = row.get("rank", "?")

    # ── 头像（左侧，垂直居中）──
    AV      = Cm(2.6)
    AV_L    = cx + RED_BAR + Cm(0.3)
    AV_T    = cy + (ch - AV) / 2

    av_buf = av_cache.get(login)
    placed = False
    if av_buf:
        try:
            av_buf.seek(0)
            sl.shapes.add_picture(av_buf, int(AV_L), int(AV_T), int(AV), int(AV))
            placed = True
        except Exception:
            pass
    if not placed:
        _rect(sl, int(AV_L), int(AV_T), int(AV), int(AV), fill=C_RED)
        _txt(sl, login[:1].upper(),
             int(AV_L), int(AV_T + AV * 0.18), int(AV), int(AV * 0.64),
             sz=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── 文字区 ──
    TL = AV_L + AV + Cm(0.3)
    TW = cx + cw - TL - Cm(0.25)
    ty = cy + Cm(0.38)

    # 排名前缀 + 姓名
    rank_pfx = {1: "🥇", 2: "🥈", 3: "🥉"}.get(
        int(rank) if str(rank).isdigit() else 0, f"#{rank}"
    )
    _txt(sl, f"{rank_pfx}  {name}",
         int(TL), int(ty), int(TW), int(Cm(0.75)),
         sz=12, bold=True, color=C_DGRAY, wrap=False)
    ty += Cm(0.78)

    # @login（红色）
    _txt(sl, f"@{login}",
         int(TL), int(ty), int(TW), int(Cm(0.6)),
         sz=10, color=C_RED, wrap=False)
    ty += Cm(0.65)

    # 💻 总提交数
    try:
        commits = int(row.get("total_commits") or 0)
        commits_str = f"{commits:,}"
    except (ValueError, TypeError):
        commits_str = "—"
    _txt(sl, "💻 总提交数",
         int(TL), int(ty), int(TW * 0.5), int(Cm(0.55)),
         sz=8, color=C_MGRAY, wrap=False)
    _txt(sl, commits_str,
         int(TL + TW * 0.5), int(ty), int(TW * 0.5), int(Cm(0.55)),
         sz=9, bold=True, color=C_RED, align=PP_ALIGN.RIGHT, wrap=False)
    ty += Cm(0.6)

    # 红色细分割线
    _rect(sl, int(TL), int(ty), int(TW * 0.92), int(Cm(0.05)), fill=C_RED)
    ty += Cm(0.18)

    # 联系/基本信息
    LINE_H = Cm(0.56)

    def _blog_url(v):
        if not v or str(v) in ("None", ""):
            return None
        return v if str(v).startswith("http") else "https://" + str(v)

    def _twitter_url(v):
        if not v or str(v) in ("None", ""):
            return None
        handle = str(v).lstrip("@")
        return f"https://twitter.com/{handle}"

    info_fields = [
        ("🔗", "GitHub", row.get("profile_url"),  row.get("profile_url")),
        ("🏢", "公司",   row.get("company"),       None),
        ("📍", "地区",   row.get("location"),      None),
        ("📧", "邮箱",   row.get("email"),          f"mailto:{row.get('email')}" if row.get("email") and str(row.get("email")) not in ("None", "") else None),
        ("🌐", "主页",   row.get("blog"),           _blog_url(row.get("blog"))),
        ("🐦", "推特",   row.get("twitter_username"), _twitter_url(row.get("twitter_username"))),
    ]
    for icon, label, val, link in info_fields:
        if ty + LINE_H > cy + ch - Cm(0.25):
            break
        if val and str(val) not in ("None", "", "0"):
            _txt(sl, f"{icon} {label}：{val}",
                 int(TL), int(ty), int(TW), int(LINE_H),
                 sz=8.5, color=C_RED if link else C_DGRAY, wrap=False,
                 url=link)
            ty += LINE_H


def _slide_summary(prs, rows: list, repo: str, pg_start: int,
                   total: int, av_cache: dict) -> int:
    """
    贡献者名片汇总页，固定 3 列 × 2 行，超出自动分页。
    返回实际生成的页数。
    """
    pages_added = 0
    batch_total = math.ceil(len(rows) / _SUMMARY_PER_PAGE) if rows else 1

    for batch_idx, batch_start in enumerate(range(0, max(len(rows), 1), _SUMMARY_PER_PAGE)):
        batch = rows[batch_start: batch_start + _SUMMARY_PER_PAGE]
        pg = pg_start + batch_idx
        pages_added += 1

        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

        sub = f"共 {len(rows)} 位贡献者  ·  精简名片"
        if batch_total > 1:
            sub += f"  （第 {batch_idx + 1} / {batch_total} 页）"
        _header(sl, "贡献者名片总览", sub)
        _footer(sl, repo, f"{pg} / {total}")

        for i, row in enumerate(batch):
            ci = i % _SUMMARY_COLS
            ri = i // _SUMMARY_COLS
            cx = _PAD_X + ci * (_CARD_W + _GAP_X)
            cy = BODY_T + _PAD_Y + ri * (_CARD_H + _GAP_Y)
            _name_card(sl, cx, cy, row, av_cache)

    return pages_added


# ════════════════════════════════════════════════════════════
# Slide N：贡献者详情
# ════════════════════════════════════════════════════════════

def _contrib_chart(row: dict, df: pd.DataFrame) -> io.BytesIO:
    top10 = df.nsmallest(10, "rank")

    def _f(v):
        try:
            return float(v or 0)
        except (TypeError, ValueError):
            return 0.0

    labels = ["提交数", "新增(k)", "删除(k)", "净增(k)"]
    person_vals = [
        _f(row.get("total_commits")),
        _f(row.get("total_additions")) / 1000,
        _f(row.get("total_deletions")) / 1000,
        max(_f(row.get("net_lines")), 0) / 1000,
    ]
    top10_vals = [
        float(top10["total_commits"].mean()),
        float(top10["total_additions"].mean()) / 1000,
        float(top10["total_deletions"].mean()) / 1000,
        float(top10["net_lines"].clip(lower=0).mean()) / 1000,
    ]

    x = np.arange(len(labels))
    w = 0.36
    fig, ax = plt.subplots(figsize=(9.5, 3.0), facecolor="white")
    b1 = ax.bar(x - w / 2, person_vals, w, label="本人",
                color="#C7000B", zorder=3)
    b2 = ax.bar(x + w / 2, top10_vals, w, label="Top10 均值",
                color="#6B0004", alpha=0.7, zorder=3)
    ax.bar_label(b1, fmt="%.1f", padding=2, fontsize=8, color="#333")
    ax.bar_label(b2, fmt="%.1f", padding=2, fontsize=8, color="#333")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, color="#6B6B6B")
    ax.legend(fontsize=9, framealpha=0, loc="upper right")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#E0E0E0")
    ax.set_facecolor("white")
    ax.grid(axis="y", linestyle="--", alpha=0.3, color="#E0E0E0", zorder=0)
    ax.set_title("贡献指标 vs Top10 均值", fontsize=10, color="#333", pad=6)
    fig.tight_layout(pad=0.5)
    return _fig_buf(fig)


def _metric_card(sl, l, t, w, h, val: str, label: str, accent=True):
    """右侧指标卡：白底 + 顶部红线（或无） + 数值 + 标签。"""
    _rect(sl, int(l), int(t), int(w), int(h),
          fill=C_WHITE, line_rgb=C_BORDER, lw_pt=0.5)
    if accent:
        _rect(sl, int(l), int(t), int(w), int(Cm(0.15)), fill=C_RED)
    _txt(sl, val,
         int(l + Cm(0.25)), int(t + Cm(0.3)), int(w - Cm(0.5)), int(h * 0.52),
         sz=15, bold=True, color=C_DGRAY, align=PP_ALIGN.CENTER)
    _txt(sl, label,
         int(l + Cm(0.15)), int(t + Cm(0.3) + int(h * 0.52)), int(w - Cm(0.3)), int(h * 0.4),
         sz=7.5, color=C_MGRAY, align=PP_ALIGN.CENTER)


def _slide_contributor(prs, row: dict, df: pd.DataFrame, pg: int, total: int,
                       av_cache: dict = None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

    login = str(row.get("login") or "unknown")
    name  = str(row.get("name") or login)
    rank  = row.get("rank", "?")

    _header(sl, "贡献者详情", f"@{login}  ·  贡献排名 #{rank}")
    _footer(sl, f"github.com/{login}", f"{pg} / {total}")

    # ════ 左侧面板 ════════════════════════════════════════════
    LW = Cm(10.5)
    # 浅灰背景
    _rect(sl, 0, BODY_T, LW, BODY_H, fill=C_LGRAY)
    # 顶部红条
    _rect(sl, 0, BODY_T, LW, Cm(0.18), fill=C_RED)

    # Avatar（优先从缓存取）
    AV_SIZE = Cm(2.8)
    AV_L, AV_T = Cm(0.7), BODY_T + Cm(0.55)
    av_buf = (av_cache or {}).get(login) or _fetch_avatar(row.get("avatar_url"))
    if av_buf:
        try:
            av_buf.seek(0)
            sl.shapes.add_picture(av_buf, int(AV_L), int(AV_T), int(AV_SIZE), int(AV_SIZE))
        except Exception:
            av_buf = None
    if not av_buf:
        _rect(sl, int(AV_L), int(AV_T), int(AV_SIZE), int(AV_SIZE), fill=C_RED)
        _txt(sl, login[:1].upper(),
             int(AV_L), int(AV_T + Cm(0.65)), int(AV_SIZE), Cm(1.5),
             sz=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 姓名 & login
    NL = AV_L + AV_SIZE + Cm(0.45)
    _txt(sl, name, int(NL), int(AV_T + Cm(0.2)), int(LW - NL - Cm(0.3)), Cm(1.1),
         sz=14, bold=True, color=C_DGRAY)
    _txt(sl, f"@{login}", int(NL), int(AV_T + Cm(1.35)), int(LW - NL - Cm(0.3)), Cm(0.65),
         sz=10, color=C_RED)

    # Bio
    bio_y = AV_T + AV_SIZE + Cm(0.4)
    if row.get("bio") and str(row["bio"]) not in ("None", ""):
        _txt(sl, str(row["bio"]), Cm(0.7), bio_y, LW - Cm(1.0), Cm(1.2),
             sz=8.5, italic=True, color=C_MGRAY, wrap=True)
        info_y = bio_y + Cm(1.35)
    else:
        info_y = bio_y

    # 红色细分割线
    _rect(sl, int(Cm(0.7)), int(info_y - Cm(0.1)),
          int(LW - Cm(1.1)), int(Cm(0.06)), fill=C_RED)

    # 联系信息
    def _blog_url(v):
        if not v or str(v) in ("None", ""):
            return None
        return v if str(v).startswith("http") else "https://" + str(v)

    def _twitter_url(v):
        if not v or str(v) in ("None", ""):
            return None
        handle = str(v).lstrip("@")
        return f"https://twitter.com/{handle}"

    fields = [
        ("🔗", "GitHub", row.get("profile_url"),      row.get("profile_url")),
        ("🏢", "公司",   row.get("company"),           None),
        ("📍", "地区",   row.get("location"),          None),
        ("📧", "邮箱",   row.get("email"),              f"mailto:{row.get('email')}" if row.get("email") and str(row.get("email")) not in ("None", "") else None),
        ("🐦", "推特",   row.get("twitter_username"),  _twitter_url(row.get("twitter_username"))),
        ("🌐", "主页",   row.get("blog"),               _blog_url(row.get("blog"))),
    ]
    y = info_y + Cm(0.15)
    for icon, label, val, link in fields:
        if val and str(val) not in ("None", "", "0"):
            _txt(sl, f"{icon} {label}：{val}", Cm(0.7), y, LW - Cm(1.0), Cm(0.62),
                 sz=8.5, color=C_RED if link else C_DGRAY, wrap=False, url=link)
            y += Cm(0.7)
            if y > BODY_B - Cm(1.5):
                break

    # 求职状态 badge
    try:
        if int(row.get("hireable") or 0) == 1:
            _rect(sl, int(Cm(0.7)), int(y + Cm(0.1)), int(Cm(3.2)), int(Cm(0.52)), fill=C_GREEN)
            _txt(sl, "✅  开放求职", Cm(0.75), y + Cm(0.13), Cm(3.1), Cm(0.44),
                 sz=8.5, bold=True, color=C_WHITE)
    except (TypeError, ValueError):
        pass

    # ════ 右侧：指标卡 + 图表 ══════════════════════════════════
    RL = LW + Cm(0.4)
    RW = SW - RL - Cm(0.4)

    metrics = [
        ("💻 总提交数",  row.get("total_commits")),
        ("➕ 新增行数",  row.get("total_additions")),
        ("➖ 删除行数",  row.get("total_deletions")),
        ("📈 净增行数",  row.get("net_lines")),
        ("👥 关注者",    row.get("followers")),
        ("📦 公开仓库",  row.get("public_repos")),
    ]
    COLS   = 3
    CW_c   = RW / COLS
    CH_c   = Cm(2.3)
    CARD_T = BODY_T + Cm(0.3)

    for i, (label, val) in enumerate(metrics):
        ci = i % COLS
        ri = i // COLS
        cx = RL + ci * CW_c
        cy = CARD_T + ri * (CH_c + Cm(0.15))
        try:
            vstr = f"{int(val):,}" if val not in (None, "None") else "—"
        except (ValueError, TypeError):
            vstr = "—"
        _metric_card(sl, cx + Cm(0.1), cy, CW_c - Cm(0.2), CH_c, vstr, label)

    # 贡献对比图
    CHART_TOP = CARD_T + 2 * (CH_c + Cm(0.15)) + Cm(0.25)
    CHART_H   = BODY_B - CHART_TOP - Cm(0.15)
    c_buf = _contrib_chart(row, df)
    sl.shapes.add_picture(c_buf, int(RL), int(CHART_TOP), int(RW), int(CHART_H))


# ════════════════════════════════════════════════════════════
# 主生成函数
# ════════════════════════════════════════════════════════════

def _add_repo_slides(prs, repo: str, df: pd.DataFrame, selected_logins: list,
                     pg_offset: int, total_pages: int, av_cache: dict) -> int:
    """
    把一个仓库的所有幻灯片（封面/概览/名片汇总/详情）添加到已有 Presentation 中。
    pg_offset: 该仓库封面的页码（1-indexed）
    返回实际添加的幻灯片数量。
    """
    sel_df   = df[df["login"].isin(selected_logins)].sort_values("rank")
    sel_rows = [r.to_dict() for _, r in sel_df.iterrows()]
    n_sel    = len(sel_rows)
    n_summary_pages = max(1, math.ceil(n_sel / _SUMMARY_PER_PAGE)) if n_sel else 1

    stats = {
        "total":   len(df),
        "commits": int(df["total_commits"].sum()) if "total_commits" in df.columns else 0,
    }

    # 封面（无页码）
    _slide_cover(prs, repo, stats)
    # 概览
    _slide_overview(prs, df, repo, pg=pg_offset + 1, total=total_pages)
    # 名片汇总
    _slide_summary(prs, sel_rows, repo,
                   pg_start=pg_offset + 2, total=total_pages, av_cache=av_cache)
    # 贡献者详情
    detail_start = pg_offset + 2 + n_summary_pages
    for i, row in enumerate(sel_rows, start=detail_start):
        _slide_contributor(prs, row, df, pg=i, total=total_pages, av_cache=av_cache)

    return 1 + 1 + n_summary_pages + n_sel


def build_ppt(repo: str, df: pd.DataFrame, selected_logins: list) -> io.BytesIO:
    """单仓库 PPT（现有接口不变）。"""
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    sel_df   = df[df["login"].isin(selected_logins)].sort_values("rank")
    sel_rows = [r.to_dict() for _, r in sel_df.iterrows()]
    n_sel    = len(sel_rows)

    n_summary_pages = max(1, math.ceil(n_sel / _SUMMARY_PER_PAGE))
    total_pages = 2 + n_summary_pages + n_sel

    av_cache: dict[str, io.BytesIO | None] = {
        str(r.get("login") or ""): _fetch_avatar(r.get("avatar_url"))
        for r in sel_rows
    }

    _add_repo_slides(prs, repo, df, selected_logins,
                     pg_offset=1, total_pages=total_pages, av_cache=av_cache)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _slide_batch_cover(prs, repos_data: list):
    """总封面幻灯片：列出所有仓库名和简要统计。"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    LP_W = Cm(14.2)
    _rect(sl, 0, 0, LP_W, SH, fill=C_DARKRED)
    _rect(sl, 0, 0, LP_W, Cm(0.5), fill=C_RED)
    _rect(sl, 0, SH - Cm(0.5), LP_W, Cm(0.5), fill=C_RED)

    _txt(sl, "批量贡献者分析报告",
         Cm(1.2), Cm(3.8), LP_W - Cm(1.8), Cm(1.3),
         sz=14, color=RGBColor(0xFF, 0xCC, 0xCC))

    total_repos  = len(repos_data)
    total_logins = sum(len(item["logins"]) for item in repos_data)

    for i, (val, label) in enumerate([
        (f"{total_repos}",        "仓库数量"),
        (f"{total_logins:,}",     "贡献者总计"),
    ]):
        y = Cm(9.5) + i * Cm(2.4)
        _txt(sl, val,   Cm(1.2), y,           LP_W - Cm(1.8), Cm(1.4),
             sz=26, bold=True, color=C_WHITE)
        _txt(sl, label, Cm(1.2), y + Cm(1.4), LP_W - Cm(1.8), Cm(0.7),
             sz=9, color=RGBColor(0xBB, 0xCC, 0xDD))

    RP_L = LP_W + Cm(0.08)
    RP_W = SW - RP_L
    _rect(sl, LP_W, 0, Cm(0.08), SH, fill=C_RED)
    _rect(sl, RP_L, 0, RP_W, SH, fill=C_WHITE)

    _txt(sl, "GitHub",
         RP_L + Cm(2), Cm(3.5), RP_W - Cm(2.5), Cm(1.5),
         sz=36, bold=True, color=C_RED)
    _txt(sl, "Contributor Analyzer",
         RP_L + Cm(2), Cm(5.1), RP_W - Cm(2.5), Cm(1.2),
         sz=20, bold=False, color=C_DGRAY)
    _txt(sl, datetime.date.today().strftime("%Y  /  %m  /  %d"),
         RP_L + Cm(2), Cm(6.5), RP_W - Cm(2.5), Cm(0.8),
         sz=11, color=C_MGRAY)

    _txt(sl, "包含仓库：",
         RP_L + Cm(2), Cm(8.0), RP_W - Cm(2.5), Cm(0.65),
         sz=10, bold=True, color=C_DGRAY)
    for i, item in enumerate(repos_data[:8]):
        _txt(sl, f"• {item['repo']}",
             RP_L + Cm(2), Cm(8.75) + i * Cm(0.78), RP_W - Cm(2.5), Cm(0.68),
             sz=10, color=C_RED)
    if len(repos_data) > 8:
        extra = len(repos_data) - 8
        _txt(sl, f"  … 及另外 {extra} 个仓库",
             RP_L + Cm(2), Cm(8.75) + 8 * Cm(0.78), RP_W - Cm(2.5), Cm(0.68),
             sz=9, color=C_MGRAY)

    _rect(sl, RP_L, SH - Cm(3.5), RP_W, Cm(3.5), fill=C_LGRAY)
    _txt(sl, "Powered by GitHub REST API",
         RP_L + Cm(2), SH - Cm(2.8), RP_W - Cm(3), Cm(0.8),
         sz=9, color=C_MGRAY)


def build_batch_ppt(repos_data: list) -> io.BytesIO:
    """
    合并多仓库 PPT。
    repos_data = [{"repo": "owner/repo", "df": df, "logins": [...]}]
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    # 预计算总页数：1（总封面）+ 每个仓库的页数
    repos_page_counts = []
    total_pages = 1  # 总封面
    for item in repos_data:
        n_sel     = len(item["logins"])
        n_summary = max(1, math.ceil(n_sel / _SUMMARY_PER_PAGE)) if n_sel else 1
        count     = 1 + 1 + n_summary + n_sel   # cover + overview + summary + detail
        repos_page_counts.append(count)
        total_pages += count

    # 预拉取所有头像（跨仓库去重）
    all_av_cache: dict[str, io.BytesIO | None] = {}
    for item in repos_data:
        sel_df = item["df"][item["df"]["login"].isin(item["logins"])]
        for _, row in sel_df.iterrows():
            login = str(row.get("login") or "")
            if login not in all_av_cache:
                all_av_cache[login] = _fetch_avatar(row.get("avatar_url"))

    # Slide 1：总封面
    _slide_batch_cover(prs, repos_data)

    # 逐仓库添加幻灯片段落
    pg_offset = 2
    for item, count in zip(repos_data, repos_page_counts):
        av_cache = {login: all_av_cache.get(login) for login in item["logins"]}
        _add_repo_slides(prs, item["repo"], item["df"], item["logins"],
                         pg_offset=pg_offset, total_pages=total_pages, av_cache=av_cache)
        pg_offset += count

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════════════════════
# 辅助函数（UI 层）
# ════════════════════════════════════════════════════════════

def _extract_region(loc: str) -> str:
    parts = str(loc).split(",")
    r = parts[-1].strip()
    if re.match(r"^\d[\d\s\-]*$", r) or r == "":
        r = parts[-2].strip() if len(parts) >= 2 else r
    return r


_MEDALS = {1: "🥇", 2: "🥈", 3: "🥉"}


def _contrib_label(row: dict) -> str:
    rank   = int(row.get("rank", 0)) if str(row.get("rank", "")).isdigit() else 0
    prefix = _MEDALS.get(rank, f"#{row.get('rank', '?')}")
    nm = f" ({row['name']})" if row.get("name") and str(row["name"]) not in ("None", "") else ""
    co = (
        f" · {str(row['company']).strip().lstrip('@')}"
        if row.get("company") and str(row["company"]) not in ("None", "") else ""
    )
    return f"{prefix} @{row['login']}{nm}{co}"


def _load_repo_df(repo: str) -> pd.DataFrame:
    raw = get_contributors(repo)
    if not raw:
        return pd.DataFrame()
    df = pd.DataFrame(raw)
    for c in ["total_commits", "total_additions", "total_deletions",
              "net_lines", "total_changes", "followers", "public_repos"]:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    return df


# ════════════════════════════════════════════════════════════
# Streamlit UI
# ════════════════════════════════════════════════════════════

st.title("📊 PPT 报告生成器")
st.caption(
    "参考华为官方浅色版模板（白底 · 华为红 #C7000B），生成 16:9 分析报告 PPT。"
    "选择 1 个仓库生成单仓库 PPT；选择多个仓库生成合并 PPT。"
)

repos = list_repos()
if not repos:
    st.warning("数据库为空，请先在「数据采集」页面爬取至少一个仓库。")
    st.stop()

all_repo_names = [r["full_name"] for r in repos]
_all_tags = list_tags()
_repo_tag_map = get_all_repo_tags()   # {repo: [{"id", "name", "color"}, ...]}

# ── 初始化 session state ──────────────────────────────────
if "ppt_repos" not in st.session_state:
    st.session_state["ppt_repos"] = []

# ════════════════════════════════════════════════════════════
# Step 1：选择仓库（双栏）
# ════════════════════════════════════════════════════════════

st.subheader("① 选择仓库")

col_repo_l, col_repo_r = st.columns(2)

with col_repo_l:
    st.markdown("**全部仓库**")

    # 标签筛选（有标签时显示）
    if _all_tags:
        tag_filter = st.multiselect(
            "按标签筛选", [t["name"] for t in _all_tags],
            key="ppt_tag_filter", placeholder="选择标签（不选=显示全部）",
            label_visibility="collapsed",
        )
        if tag_filter:
            _filter_tag_ids = [t["id"] for t in _all_tags if t["name"] in tag_filter]
            _tagged_repos = set(get_repos_by_tags(_filter_tag_ids))
            _base_repos = [r for r in all_repo_names if r in _tagged_repos]
        else:
            _base_repos = all_repo_names
    else:
        _base_repos = all_repo_names

    repo_search = st.text_input(
        "搜索仓库", placeholder="输入关键词筛选...",
        key="ppt_repo_search", label_visibility="collapsed",
    )
    kw = repo_search.strip().lower()
    visible_repos = [r for r in _base_repos if kw in r.lower()] if kw else _base_repos

    # ── 全选 / 取消全选（操作当前可见列表）──
    def _select_all_visible(vlist=None):
        for r in (vlist or []):
            st.session_state[f"ppt_cb_{r}"] = True
            if r not in st.session_state["ppt_repos"]:
                st.session_state["ppt_repos"].append(r)
            if r not in st.session_state["ppt_contribs"]:
                st.session_state["ppt_contribs"][r] = []

    def _deselect_all_visible(vlist=None):
        for r in (vlist or []):
            st.session_state[f"ppt_cb_{r}"] = False
            if r in st.session_state["ppt_repos"]:
                st.session_state["ppt_repos"].remove(r)

    import functools
    sa_col, da_col = st.columns(2)
    with sa_col:
        st.button(
            "全选" if not kw else f"全选结果（{len(visible_repos)}）",
            key="ppt_select_all",
            use_container_width=True,
            on_click=functools.partial(_select_all_visible, visible_repos),
        )
    with da_col:
        st.button(
            "取消全选",
            key="ppt_deselect_all",
            use_container_width=True,
            on_click=functools.partial(_deselect_all_visible, visible_repos),
        )

    def _tag_badges_html(repo_name: str, font_size: str = "0.75rem") -> str:
        tags = _repo_tag_map.get(repo_name, [])
        if not tags:
            return ""
        return " ".join(
            f"<span style='background:{t['color']};color:#fff;border-radius:3px;"
            f"padding:1px 6px;font-size:{font_size};display:inline-block;margin:1px'>"
            f"{t['name']}</span>"
            for t in tags
        )

    for repo in visible_repos:
        checked = st.checkbox(f"📦 {repo}", key=f"ppt_cb_{repo}")
        if checked and repo not in st.session_state["ppt_repos"]:
            st.session_state["ppt_repos"].append(repo)
            if repo not in st.session_state["ppt_contribs"]:
                st.session_state["ppt_contribs"][repo] = []
        elif not checked and repo in st.session_state["ppt_repos"]:
            st.session_state["ppt_repos"].remove(repo)
        badges_html = _tag_badges_html(repo)
        if badges_html:
            st.markdown(
                f"<div style='margin:-8px 0 4px 26px'>{badges_html}</div>",
                unsafe_allow_html=True,
            )

with col_repo_r:
    st.markdown("**已选仓库**")
    if not st.session_state["ppt_repos"]:
        st.caption("← 从左侧勾选仓库")
    else:
        for repo in list(st.session_state["ppt_repos"]):
            rc1, rc2 = st.columns([6, 1])
            with rc1:
                n_c = len(st.session_state["ppt_contribs"].get(repo, []))
                st.markdown(f"📦 `{repo}`" + (f"  · **{n_c} 人已选**" if n_c else ""))
                badges_html = _tag_badges_html(repo)
                if badges_html:
                    st.markdown(badges_html, unsafe_allow_html=True)
            with rc2:
                if st.button("×", key=f"ppt_rm_repo_{repo}", help="移除该仓库"):
                    st.session_state["ppt_repos"].remove(repo)
                    st.session_state[f"ppt_cb_{repo}"] = False
                    st.rerun()

# ════════════════════════════════════════════════════════════
# Step 2：筛选 · 预览 · 导出 / 生成 PPT
# ════════════════════════════════════════════════════════════

if st.session_state["ppt_repos"]:
    st.markdown("---")
    st.subheader("② 筛选贡献者")

    # ── 合并所有选定仓库的贡献者数据（带缓存）────────────────
    @st.cache_data(ttl=60, show_spinner=False)
    def _load_combined_ppt(repo_tuple: tuple) -> pd.DataFrame:
        frames = []
        for rname in repo_tuple:
            df_r = _load_repo_df(rname)
            if not df_r.empty:
                df_r = df_r.copy()
                df_r["_repo"] = rname
                frames.append(df_r)
        if not frames:
            return pd.DataFrame()
        combined = pd.concat(frames, ignore_index=True)
        for col in ("company", "location"):
            if col in combined.columns:
                combined[col] = combined[col].fillna("").str.strip()
        return combined

    df_all = _load_combined_ppt(tuple(st.session_state["ppt_repos"]))

    if df_all.empty:
        st.warning("所选仓库暂无贡献者数据。")
        st.stop()

    # ── 公司 / 地区 dropdown ─────────────────────────────────
    def _clean_opts(series: pd.Series) -> list:
        return sorted(
            {v.lstrip("@").strip() for v in series if v and v.strip()},
            key=str.lower,
        )

    company_opts  = _clean_opts(df_all["company"])  if "company"  in df_all.columns else []
    location_opts = _clean_opts(df_all["location"]) if "location" in df_all.columns else []

    fc1, fc2 = st.columns(2)
    with fc1:
        sel_cos = st.multiselect(
            "🏢 公司", company_opts,
            placeholder="选择公司（可多选，不选=全部）…",
            key="ppt_cs_co",
        )
    with fc2:
        sel_locs = st.multiselect(
            "📍 地区", location_opts,
            placeholder="选择地区（可多选，不选=全部）…",
            key="ppt_cs_loc",
        )

    # ── 过滤逻辑（两项同时填写取交集）──────────────────────
    df_filtered = df_all.copy()
    if sel_cos:
        df_filtered = df_filtered[
            df_filtered["company"].str.lstrip("@").str.strip().isin(sel_cos)
        ]
    if sel_locs:
        df_filtered = df_filtered[df_filtered["location"].isin(sel_locs)]
    df_filtered = df_filtered.sort_values(["_repo", "rank"]).reset_index(drop=True)

    # ── 结果表格 ─────────────────────────────────────────────
    st.markdown("---")
    n_repos_hit = df_filtered["_repo"].nunique()
    st.markdown(f"**结果：{len(df_filtered)} 位贡献者，来自 {n_repos_hit} 个仓库**")

    _display_cols = [c for c in [
        "_repo", "rank", "login", "name", "company", "location",
        "total_commits", "total_additions", "total_deletions",
        "followers", "email", "blog", "profile_url",
    ] if c in df_filtered.columns]

    st.dataframe(
        df_filtered[_display_cols],
        use_container_width=True,
        hide_index=True,
        height=420,
        column_config={
            "_repo":           st.column_config.TextColumn("仓库"),
            "rank":            st.column_config.NumberColumn("排名", width="small"),
            "login":           st.column_config.TextColumn("用户名"),
            "name":            st.column_config.TextColumn("姓名"),
            "company":         st.column_config.TextColumn("公司"),
            "location":        st.column_config.TextColumn("地区"),
            "total_commits":   st.column_config.NumberColumn("Commits",  format="%d"),
            "total_additions": st.column_config.NumberColumn("新增行",   format="%d"),
            "total_deletions": st.column_config.NumberColumn("删除行",   format="%d"),
            "followers":       st.column_config.NumberColumn("Followers", format="%d"),
            "email":           st.column_config.TextColumn("邮箱"),
            "blog":            st.column_config.LinkColumn("主页"),
            "profile_url":     st.column_config.LinkColumn("GitHub"),
        },
    )

    if df_filtered.empty:
        st.info("没有符合条件的贡献者，请调整筛选条件。")
        st.stop()

    # ── 操作区：导出 CSV / 生成 PPT ──────────────────────────
    st.markdown("---")
    ac1, ac2 = st.columns(2)

    with ac1:
        csv_bytes = (
            df_filtered[_display_cols]
            .rename(columns={"_repo": "仓库"})
            .to_csv(index=False, encoding="utf-8-sig")
            .encode("utf-8-sig")
        )
        st.download_button(
            "⬇️ 导出结果 CSV",
            data=csv_bytes,
            file_name="filtered_contributors.csv",
            mime="text/csv",
            use_container_width=True,
        )

    with ac2:
        # 预估页数
        repos_logins: dict[str, list] = {}
        for _, row in df_filtered.iterrows():
            repos_logins.setdefault(row["_repo"], []).append(row["login"])
        n_multi = len(repos_logins) > 1
        total_pages = 1 if n_multi else 0
        for logins in repos_logins.values():
            n_s = len(logins)
            n_sum = max(1, math.ceil(n_s / _SUMMARY_PER_PAGE)) if n_s else 1
            total_pages += 1 + 1 + n_sum + n_s

        if st.button(
            f"🚀 生成 PPT（预计 {total_pages} 页）",
            type="primary",
            use_container_width=True,
            key="ppt_gen_btn",
        ):
            with st.spinner("正在生成 PPT，拉取头像约需数秒…"):
                try:
                    repos_data = []
                    for repo, logins in repos_logins.items():
                        rdf = _load_repo_df(repo)
                        if rdf.empty:
                            st.warning(f"⚠️ {repo} 暂无数据，已跳过")
                            continue
                        repos_data.append({"repo": repo, "df": rdf, "logins": logins})

                    if not repos_data:
                        st.error("没有可用数据，请检查仓库。")
                        st.stop()

                    if len(repos_data) == 1:
                        item = repos_data[0]
                        ppt_buf = build_ppt(item["repo"], item["df"], item["logins"])
                        fname = f"contributors_{item['repo'].replace('/', '_')}.pptx"
                    else:
                        ppt_buf = build_batch_ppt(repos_data)
                        slug = "_".join(r["repo"].replace("/", "-") for r in repos_data[:3])
                        fname = f"batch_contributors_{slug}.pptx"

                    st.download_button(
                        label="⬇️ 点击下载 PPT（可直接用 Google Slides 打开）",
                        data=ppt_buf,
                        file_name=fname,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True,
                    )
                    st.success(f"✅ PPT 生成成功！共 {total_pages} 页。")
                except Exception as e:
                    st.error(f"生成失败：{e}")
                    st.code(traceback.format_exc())
