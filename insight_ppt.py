"""
洞察报告 PPT 构建模块（自含原语，不 import Streamlit）。

用法：
    from insight_ppt import build_insight_ppt, THEMES
    buf = build_insight_ppt(regions_talents, repos_info, llm_content, theme, report_title)
"""

import io
import datetime
import math
import urllib.request

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 主题系统 ─────────────────────────────────────────────────
THEMES = {
    "华为经典": {
        "primary":    (0xC7, 0x00, 0x0B),
        "primary_mid":(0xA3, 0x00, 0x09),
        "bg_panel":   (0x6B, 0x00, 0x04),
        "text_main":  (0x33, 0x33, 0x33),
        "text_sub":   (0x6B, 0x6B, 0x6B),
        "bg_card":    (0xF5, 0xF5, 0xF5),
        "border":     (0xE0, 0xE0, 0xE0),
        "swatch":     "#C7000B",
    },
    "深海蓝": {
        "primary":    (0x1A, 0x56, 0xAB),
        "primary_mid":(0x2E, 0x6D, 0xCF),
        "bg_panel":   (0x0D, 0x2E, 0x5E),
        "text_main":  (0x1A, 0x24, 0x3A),
        "text_sub":   (0x55, 0x68, 0x84),
        "bg_card":    (0xF0, 0xF5, 0xFF),
        "border":     (0xC8, 0xD8, 0xF0),
        "swatch":     "#1A56AB",
    },
    "森林绿": {
        "primary":    (0x1E, 0x7A, 0x4E),
        "primary_mid":(0x2D, 0xA0, 0x69),
        "bg_panel":   (0x0D, 0x3D, 0x27),
        "text_main":  (0x1A, 0x2E, 0x22),
        "text_sub":   (0x52, 0x73, 0x5E),
        "bg_card":    (0xF0, 0xFA, 0xF4),
        "border":     (0xC2, 0xE5, 0xD0),
        "swatch":     "#1E7A4E",
    },
    "暮光紫": {
        "primary":    (0x6B, 0x3F, 0xA0),
        "primary_mid":(0x8B, 0x5C, 0xC8),
        "bg_panel":   (0x35, 0x1A, 0x55),
        "text_main":  (0x22, 0x1A, 0x33),
        "text_sub":   (0x6A, 0x57, 0x84),
        "bg_card":    (0xF8, 0xF3, 0xFF),
        "border":     (0xD9, 0xC8, 0xF0),
        "swatch":     "#6B3FA0",
    },
    "极简灰": {
        "primary":    (0x2D, 0x3A, 0x4A),
        "primary_mid":(0x4A, 0x5C, 0x72),
        "bg_panel":   (0x16, 0x1C, 0x25),
        "text_main":  (0x1A, 0x1A, 0x1A),
        "text_sub":   (0x66, 0x66, 0x66),
        "bg_card":    (0xF4, 0xF5, 0xF6),
        "border":     (0xD8, 0xDA, 0xDE),
        "swatch":     "#2D3A4A",
    },
}

# ── 当前活跃颜色（由 _apply_theme() 写入）─────────────────────
C_RED    = RGBColor(0xC7, 0x00, 0x0B)
C_DRED   = RGBColor(0xA3, 0x00, 0x09)
C_DARKRED= RGBColor(0x6B, 0x00, 0x04)
C_DGRAY  = RGBColor(0x33, 0x33, 0x33)
C_MGRAY  = RGBColor(0x6B, 0x6B, 0x6B)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ── 尺寸常量 ──────────────────────────────────────────────────
SW     = Cm(33.87)   # 16:9 宽
SH     = Cm(19.05)   # 16:9 高
HDR_H  = Cm(2.6)     # 页眉高
FTR_H  = Cm(0.9)     # 页脚高
BODY_T = HDR_H       # 内容起始 Y
BODY_B = SH - FTR_H  # 内容结束 Y
BODY_H = BODY_B - BODY_T


def _apply_theme(theme_name: str):
    """将主题色系应用到全局 C_* 变量。需在 build 前调用。"""
    global C_RED, C_DRED, C_DARKRED, C_DGRAY, C_MGRAY, C_LGRAY, C_BORDER
    t = THEMES.get(theme_name, THEMES["华为经典"])
    C_RED     = RGBColor(*t["primary"])
    C_DRED    = RGBColor(*t["primary_mid"])
    C_DARKRED = RGBColor(*t["bg_panel"])
    C_DGRAY   = RGBColor(*t["text_main"])
    C_MGRAY   = RGBColor(*t["text_sub"])
    C_LGRAY   = RGBColor(*t["bg_card"])
    C_BORDER  = RGBColor(*t["border"])


# ════════════════════════════════════════════════════════════
# 低级绘图原语
# ════════════════════════════════════════════════════════════

def _rect(sl, l, t, w, h, fill=None, line_rgb=None, lw_pt=0.5):
    """添加矩形色块。"""
    s = sl.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(lw_pt)
    else:
        s.line.fill.background()
    return s


def _txt(sl, text, l, t, w, h,
         sz=11, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, wrap=True, url=None):
    """添加文字框。color=None 时使用当前主题深灰色。"""
    if color is None:
        color = C_DGRAY
    tb = sl.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text) if text not in (None, "None") else ""
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Arial"
    r.font.color.rgb = color
    if url:
        r.hyperlink.address = url
        r.font.underline = True
    return tb


def _fetch_avatar(url: str):
    """下载头像图片，返回 BytesIO 或 None。"""
    if not url:
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "PPT-Generator/1.0"})
        with urllib.request.urlopen(req, timeout=5) as r:
            buf = io.BytesIO(r.read())
            buf.seek(0)
            return buf
    except Exception:
        return None


# ════════════════════════════════════════════════════════════
# 通用页眉 / 页脚
# ════════════════════════════════════════════════════════════

def _header(sl, title: str, sub: str = ""):
    _rect(sl, 0, 0, SW, HDR_H, fill=C_WHITE)
    _rect(sl, 0, 0, Cm(0.35), HDR_H, fill=C_RED)
    _rect(sl, 0, HDR_H - Cm(0.08), SW, Cm(0.08), fill=C_BORDER)
    _txt(sl, title, Cm(0.75), Cm(0.4), SW - Cm(2), Cm(1.3),
         sz=20, bold=True, color=C_DGRAY)
    if sub:
        _txt(sl, sub, Cm(0.75), Cm(1.75), SW - Cm(2), Cm(0.72),
             sz=10, color=C_MGRAY)


def _blank_layout(prs):
    """Return a blank-ish slide layout safely (layout[6] or fallback to last)."""
    layouts = prs.slide_layouts
    idx = min(6, len(layouts) - 1)
    return layouts[idx]


def _footer(sl, note: str = ""):
    _rect(sl, 0, BODY_B, SW, FTR_H, fill=C_LGRAY)
    _rect(sl, 0, BODY_B, SW, Cm(0.08), fill=C_RED)
    if note:
        _txt(sl, note, SW - Cm(5.5), BODY_B + Cm(0.15), Cm(5.2), FTR_H - Cm(0.2),
             sz=8, color=C_MGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# Slide 1：封面
# ════════════════════════════════════════════════════════════

def slide_cover(prs, title: str, subtitle: str, date: str):
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

    # ── 左侧深色面板 ──
    LP_W = Cm(14.2)
    _rect(sl, 0, 0, LP_W, SH, fill=C_DARKRED)
    _rect(sl, 0, 0, LP_W, Cm(0.5), fill=C_RED)
    _rect(sl, 0, SH - Cm(0.5), LP_W, Cm(0.5), fill=C_RED)

    _txt(sl, "开源人才洞察报告",
         Cm(1.2), Cm(3.8), LP_W - Cm(1.8), Cm(1.3),
         sz=14, color=RGBColor(0xFF, 0xCC, 0xCC))
    _txt(sl, title,
         Cm(1.2), Cm(5.2), LP_W - Cm(1.8), Cm(3.5),
         sz=22, bold=True, color=C_WHITE, wrap=True)
    _rect(sl, Cm(1.2), Cm(9.0), Cm(5), Cm(0.06),
          fill=RGBColor(0x88, 0x99, 0xBB))
    _txt(sl, date,
         Cm(1.2), Cm(9.3), LP_W - Cm(1.8), Cm(0.8),
         sz=10, color=RGBColor(0xBB, 0xCC, 0xDD))

    # ── 右侧白色内容区 ──
    RP_L = LP_W + Cm(0.08)
    RP_W = SW - RP_L
    _rect(sl, LP_W, 0, Cm(0.08), SH, fill=C_RED)
    _rect(sl, RP_L, 0, RP_W, SH, fill=C_WHITE)

    _txt(sl, title,
         RP_L + Cm(2), Cm(3.5), RP_W - Cm(2.5), Cm(2.8),
         sz=28, bold=True, color=C_RED, wrap=True)
    if subtitle:
        _txt(sl, subtitle,
             RP_L + Cm(2), Cm(6.5), RP_W - Cm(2.5), Cm(1.0),
             sz=14, color=C_DGRAY)
    _txt(sl, date,
         RP_L + Cm(2), Cm(7.8), RP_W - Cm(2.5), Cm(0.8),
         sz=11, color=C_MGRAY)

    _rect(sl, RP_L, SH - Cm(3.5), RP_W, Cm(3.5), fill=C_LGRAY)
    _txt(sl, "Powered by GitHub REST API  ·  AI-Enhanced",
         RP_L + Cm(2), SH - Cm(2.8), RP_W - Cm(3), Cm(0.8),
         sz=9, color=C_MGRAY)
    _txt(sl, "开源人才洞察报告",
         RP_L + Cm(2), SH - Cm(1.9), RP_W - Cm(3), Cm(0.8),
         sz=10, bold=True, color=C_RED)


# ════════════════════════════════════════════════════════════
# Slide 2：总览 - 人才质量与密度
# ════════════════════════════════════════════════════════════

def slide_overview_quality(prs, quality_summary: str, density_stats: str,
                            talent_count_by_region: dict, pg: int, total: int,
                            match_distribution: dict = None,
                            total_count: int = 0):
    """
    match_distribution : {"A": n, "B": n, "C": n, "D": n} — 各匹配级别人数
    total_count        : 全部人才总数
    """
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)
    _header(sl, "人才总览", "质量评估与地区密度分析")
    _footer(sl, f"{pg} / {total}")

    md = match_distribution or {}
    n_regions = len(talent_count_by_region)
    n_total   = total_count or sum(talent_count_by_region.values())
    n_high    = md.get("A", 0) + md.get("B", 0)  # A+B 视为高匹配

    # ── 顶部三格指标卡 ────────────────────────────────────────
    METRIC_T = BODY_T + Cm(0.3)
    METRIC_H = Cm(2.0)
    CARD_GAP  = Cm(0.3)
    usable_w  = SW - Cm(1.0)
    METRIC_W  = (usable_w - CARD_GAP * 2) / 3

    metrics = [
        ("覆盖地区", str(n_regions), "个"),
        ("收录人才", str(n_total), "人"),
        ("高匹配人才（A/B 级）", str(n_high), "人"),
    ]
    for idx, (label, value, unit) in enumerate(metrics):
        mx = Cm(0.5) + idx * (METRIC_W + CARD_GAP)
        _rect(sl, mx, METRIC_T, METRIC_W, METRIC_H,
              fill=C_LGRAY, line_rgb=C_BORDER, lw_pt=0.4)
        _rect(sl, mx, METRIC_T, METRIC_W, Cm(0.18), fill=C_RED)
        _txt(sl, label,
             mx + Cm(0.3), METRIC_T + Cm(0.25), METRIC_W - Cm(0.4), Cm(0.5),
             sz=8.5, color=C_MGRAY)
        _txt(sl, value,
             mx + Cm(0.3), METRIC_T + Cm(0.78), METRIC_W - Cm(1.5), Cm(0.85),
             sz=26, bold=True, color=C_RED)
        _txt(sl, unit,
             mx + METRIC_W - Cm(1.4), METRIC_T + Cm(1.3), Cm(1.1), Cm(0.55),
             sz=9, color=C_MGRAY)

    # ── 下方两栏 ──────────────────────────────────────────────
    col_t = METRIC_T + METRIC_H + Cm(0.3)
    col_h = BODY_B - col_t - Cm(0.2)

    # LEFT：地区密度
    COL_L_W = Cm(10.5)
    col_l = Cm(0.5)
    _rect(sl, col_l, col_t, COL_L_W, col_h, fill=C_LGRAY)

    _txt(sl, "地区人才密度",
         col_l + Cm(0.4), col_t + Cm(0.25), COL_L_W - Cm(0.6), Cm(0.55),
         sz=10, bold=True, color=C_DGRAY)
    if density_stats:
        _txt(sl, density_stats,
             col_l + Cm(0.4), col_t + Cm(0.85), COL_L_W - Cm(0.6), Cm(0.7),
             sz=8.5, color=C_MGRAY, wrap=True)

    card_y = col_t + Cm(1.65)
    max_count = max(talent_count_by_region.values(), default=1)
    for region, count in list(talent_count_by_region.items())[:7]:
        if card_y + Cm(0.95) > col_t + col_h - Cm(0.2):
            break
        CARD_H = Cm(0.85)
        _rect(sl, col_l + Cm(0.3), card_y, COL_L_W - Cm(0.5), CARD_H,
              fill=C_WHITE, line_rgb=C_BORDER, lw_pt=0.4)
        _rect(sl, col_l + Cm(0.3), card_y, Cm(0.18), CARD_H, fill=C_RED)
        # 进度条背景
        bar_l  = col_l + Cm(0.65)
        bar_t  = card_y + CARD_H - Cm(0.22)
        bar_w  = COL_L_W - Cm(2.8)
        _rect(sl, bar_l, bar_t, bar_w, Cm(0.1), fill=C_BORDER)
        _rect(sl, bar_l, bar_t, bar_w * count / max_count, Cm(0.1), fill=C_RED)
        _txt(sl, region,
             col_l + Cm(0.65), card_y + Cm(0.1), COL_L_W - Cm(2.5), Cm(0.5),
             sz=9, color=C_DGRAY, wrap=False)
        _txt(sl, f"{count} 人",
             col_l + COL_L_W - Cm(1.7), card_y + Cm(0.1), Cm(1.4), Cm(0.5),
             sz=9, bold=True, color=C_RED, align=PP_ALIGN.RIGHT, wrap=False)
        card_y += CARD_H + Cm(0.12)

    # RIGHT：匹配度分布 + 质量评估
    col_r_l = col_l + COL_L_W + Cm(0.4)
    col_r_w = SW - col_r_l - Cm(0.5)
    cur_y = col_t + Cm(0.25)

    # 匹配度分布（若有数据）
    if md:
        _txt(sl, "匹配度分布",
             col_r_l, cur_y, col_r_w, Cm(0.55),
             sz=10, bold=True, color=C_DGRAY)
        cur_y += Cm(0.6)

        level_colors = {
            "A": C_RED,
            "B": C_DRED,
            "C": C_MGRAY,
            "D": C_BORDER,
        }
        level_labels = {"A": "A 级（优先推荐）", "B": "B 级（重点关注）",
                        "C": "C 级（候选储备）", "D": "D 级（待观察）"}
        total_md = sum(md.values()) or 1
        for level in ("A", "B", "C", "D"):
            cnt = md.get(level, 0)
            if cnt == 0:
                continue
            if cur_y + Cm(0.7) > col_t + col_h * 0.48:
                break
            bar_fill_w = (col_r_w - Cm(5.5)) * cnt / total_md
            _txt(sl, level_labels.get(level, level),
                 col_r_l, cur_y, col_r_w - Cm(2.2), Cm(0.42),
                 sz=8, color=C_DGRAY)
            _txt(sl, f"{cnt} 人",
                 col_r_l + col_r_w - Cm(1.8), cur_y, Cm(1.5), Cm(0.42),
                 sz=8, bold=True, color=level_colors.get(level, C_MGRAY),
                 align=PP_ALIGN.RIGHT)
            bar_t2 = cur_y + Cm(0.44)
            _rect(sl, col_r_l, bar_t2, col_r_w - Cm(0.1), Cm(0.2), fill=C_LGRAY)
            _rect(sl, col_r_l, bar_t2, max(Cm(0.15), bar_fill_w + Cm(0.1)), Cm(0.2),
                  fill=level_colors.get(level, C_MGRAY))
            cur_y += Cm(0.72)

        cur_y += Cm(0.2)
        _rect(sl, col_r_l, cur_y, col_r_w, Cm(0.05), fill=C_BORDER)
        cur_y += Cm(0.25)

    # 质量总评
    _txt(sl, "人才质量评估",
         col_r_l, cur_y, col_r_w, Cm(0.55),
         sz=10, bold=True, color=C_DGRAY)
    cur_y += Cm(0.6)
    _rect(sl, col_r_l, cur_y, Cm(0.18), min(Cm(3.5), BODY_B - cur_y - Cm(0.3)),
          fill=C_RED)
    _txt(sl, quality_summary or "—",
         col_r_l + Cm(0.35), cur_y, col_r_w - Cm(0.4),
         max(Cm(1.0), BODY_B - cur_y - Cm(0.3)),
         sz=10, wrap=True, color=C_DGRAY)


# ════════════════════════════════════════════════════════════
# Slide 3：总览 - 华为匹配（PLACEHOLDER）
# ════════════════════════════════════════════════════════════

def slide_overview_huawei_placeholder(prs, repos_info: dict, pg: int, total: int,
                                       page_num: int = 1, total_pages_this: int = 1):
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)
    header_sub = f"项目技术价值评估  ·  第 {page_num}/{total_pages_this} 页" if total_pages_this > 1 else "项目技术价值评估"
    _header(sl, "华为匹配度分析", header_sub)
    _footer(sl, f"{pg} / {total}")

    # 项目卡片（2列，动态行数，每页最多6个）
    repos = list(repos_info.items())
    COLS = 2
    CARD_W = (SW - Cm(2.4)) / COLS
    CARD_PAD_T = BODY_T + Cm(0.4)
    AVAILABLE_H = BODY_B - CARD_PAD_T
    GAP = Cm(0.3)

    n_rows = max(1, math.ceil(len(repos) / COLS))
    CARD_H = min(Cm(4.5), (AVAILABLE_H - GAP * (n_rows - 1)) / n_rows)

    for i, (repo, info) in enumerate(repos):
        col = i % COLS
        row = i // COLS
        cx = Cm(1.0) + col * (CARD_W + Cm(0.4))
        cy = CARD_PAD_T + row * (CARD_H + GAP)

        _rect(sl, cx, cy, CARD_W, CARD_H, fill=C_LGRAY, line_rgb=C_BORDER, lw_pt=0.5)
        _rect(sl, cx, cy, Cm(0.25), CARD_H, fill=C_RED)

        short = repo.split("/")[-1] if "/" in repo else repo
        _txt(sl, short,
             cx + Cm(0.5), cy + Cm(0.25), CARD_W - Cm(0.7), Cm(0.7),
             sz=12, bold=True, color=C_DGRAY)

        # tech_areas badges
        tech_areas = info.get("tech_areas") or []
        bx = cx + Cm(0.5)
        by = cy + Cm(1.05)
        for tag in tech_areas[:3]:
            tag_w = max(Cm(1.5), Cm(len(tag) * 0.22 + 0.4))
            _rect(sl, bx, by, tag_w, Cm(0.45), fill=C_RED)
            _txt(sl, tag,
                 bx + Cm(0.1), by + Cm(0.02), tag_w - Cm(0.15), Cm(0.42),
                 sz=7.5, bold=True, color=C_WHITE, wrap=False)
            bx += tag_w + Cm(0.15)

        desc_h = CARD_H * 0.27
        desc = info.get("description") or ""
        _txt(sl, desc[:100],
             cx + Cm(0.5), cy + Cm(1.65), CARD_W - Cm(0.7), desc_h,
             sz=8.5, color=C_DGRAY, wrap=True)

        hw_val = info.get("huawei_value") or ""
        _txt(sl, hw_val[:120],
             cx + Cm(0.5), cy + Cm(1.65) + desc_h + Cm(0.1), CARD_W - Cm(0.7),
             CARD_H - Cm(1.65) - desc_h - Cm(0.2),
             sz=8, color=C_MGRAY, italic=True, wrap=True)


# ════════════════════════════════════════════════════════════
# Slide 4：地区分节页（无 header/footer chrome）
# ════════════════════════════════════════════════════════════

def slide_region_divider(prs, region_name: str, talent_count: int, pg: int, total: int):
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

    LP_W = Cm(8)
    _rect(sl, 0, 0, LP_W, SH, fill=C_DARKRED)
    _rect(sl, 0, 0, LP_W, Cm(0.5), fill=C_RED)
    _rect(sl, 0, SH - Cm(0.5), LP_W, Cm(0.5), fill=C_RED)

    # 左：大数字 + "位人才"
    _txt(sl, str(talent_count),
         Cm(1.0), Cm(5.5), LP_W - Cm(1.5), Cm(4.5),
         sz=64, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)
    _txt(sl, "位人才",
         Cm(1.0), Cm(11.0), LP_W - Cm(1.5), Cm(1.2),
         sz=18, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    # 红色竖分割线
    _rect(sl, LP_W, Cm(2), Cm(0.12), SH - Cm(4), fill=C_RED)

    # 右：地区名 + 副标题
    RP_L = LP_W + Cm(0.5)
    RP_W = SW - RP_L - Cm(1)
    _txt(sl, region_name,
         RP_L, Cm(6.0), RP_W, Cm(3.0),
         sz=36, bold=True, color=C_DGRAY, wrap=False)
    _txt(sl, "人才 · 专题分析",
         RP_L, Cm(9.3), RP_W, Cm(1.2),
         sz=16, color=C_MGRAY)

    # 页码（右下角）
    _txt(sl, f"{pg} / {total}",
         SW - Cm(4), SH - Cm(1.2), Cm(3.5), Cm(0.8),
         sz=9, color=C_MGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 溢出安全游标（面板内部使用）
# ════════════════════════════════════════════════════════════

class _PanelWriter:
    """
    在固定矩形区域内按行绘制内容，超出底部时自动跳过。
    若剩余空间不足以容纳所有 info 行，自动压缩行高/字号。
    """
    def __init__(self, sl, left, top, width, bottom, pad_l=Cm(0.4), pad_r=Cm(0.4)):
        self.sl      = sl
        self.left    = left
        self.top     = top
        self.width   = width
        self.bottom  = bottom
        self.pad_l   = pad_l
        self.pad_r   = pad_r
        self.y       = top

    @property
    def remaining(self):
        return self.bottom - self.y

    def fits(self, h):
        return self.y + h <= self.bottom

    def txt(self, text, h, sz=10, bold=False, italic=False,
            color=None, align=PP_ALIGN.LEFT, url=None, wrap=False):
        if not self.fits(h):
            return False
        _txt(self.sl, text,
             self.left + self.pad_l,
             self.y,
             self.width - self.pad_l - self.pad_r,
             h,
             sz=sz, bold=bold, italic=italic,
             color=color or C_DGRAY,
             align=align, wrap=wrap, url=url)
        self.y += h
        return True

    def hrule(self, h=Cm(0.05), gap_after=Cm(0.12), fill=None):
        if not self.fits(h + gap_after):
            return False
        _rect(self.sl,
              self.left + self.pad_l,
              self.y,
              self.width - self.pad_l - self.pad_r,
              h,
              fill=fill or C_BORDER)
        self.y += h + gap_after
        return True

    def skip(self, h):
        self.y += h

    def bar(self, fraction, h=Cm(0.3), gap_after=Cm(0.08)):
        """水平进度条，fraction ∈ [0, 1]。"""
        if not self.fits(h + gap_after):
            return False
        bar_w = self.width - self.pad_l - self.pad_r
        _rect(self.sl, self.left + self.pad_l, self.y, bar_w, h, fill=C_BORDER)
        if fraction > 0:
            _rect(self.sl, self.left + self.pad_l, self.y,
                  bar_w * min(1.0, fraction), h, fill=C_RED)
        self.y += h + gap_after
        return True

    def badge(self, text, y_offset=Cm(0.05), h=Cm(0.5), color=None):
        """在右上角画徽章（不移动 y 游标）。"""
        bw = Cm(1.5)
        bx = self.left + self.width - self.pad_r - bw
        _rect(self.sl, bx, self.y + y_offset, bw, h,
              fill=color or C_RED)
        _txt(self.sl, text,
             bx + Cm(0.08), self.y + y_offset + Cm(0.05),
             bw - Cm(0.1), h - Cm(0.08),
             sz=8.5, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, wrap=False)


# ════════════════════════════════════════════════════════════
# Slide 5：人才详情页（三区布局）
# ════════════════════════════════════════════════════════════

def _clean(val, maxlen=32):
    """清洗字段：空/None/字面量None → ''，超长截断。"""
    if not val or str(val).strip() in ("", "None", "none"):
        return ""
    return str(val).strip()[:maxlen]


def slide_talent(prs, talent_data: dict, repos: list, pg: int, total: int, av_cache: dict):
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

    # ── 字段提取 ─────────────────────────────────────────────
    login    = _clean(talent_data.get("login") or talent_data.get("username"), 40) or "?"
    name     = _clean(talent_data.get("name") or talent_data.get("fullname"), 40) or login
    region   = _clean(talent_data.get("location"), 35)
    company  = _clean(talent_data.get("company") or talent_data.get("employer"), 35)
    email    = _clean(talent_data.get("email"), 40)
    blog     = _clean(talent_data.get("blog") or talent_data.get("website"), 45)
    twitter  = _clean(talent_data.get("twitter_username"), 30)
    linkedin = _clean(talent_data.get("linkedin"), 45)
    profile_url = _clean(talent_data.get("profile_url"), 80)

    commits   = talent_data.get("total_commits") or 0
    followers = talent_data.get("followers") or talent_data.get("num_followers") or 0

    # AI 字段
    ai = talent_data.get("_ai") or {}
    tech_direction       = _clean(ai.get("tech_direction"), 120)
    contribution_summary = _clean(ai.get("contribution_summary"), 200)
    key_skills           = [s for s in (ai.get("key_skills") or []) if s]
    match_score          = ai.get("match_score")
    match_level          = _clean(ai.get("match_level"), 4)
    match_reason         = _clean(ai.get("match_reason"), 250)

    _header(sl, f"{name}  ·  人才档案", f"@{login}")
    _footer(sl, f"{pg} / {total}")

    # ── 布局常量 ──────────────────────────────────────────────
    BODY_T2    = BODY_T + Cm(0.3)
    PANEL_BTM  = BODY_B - Cm(0.12)      # 面板公共底部
    PANEL_H    = PANEL_BTM - BODY_T2

    LEFT_W      = Cm(9.5)
    left_l      = Cm(0.5)
    RIGHT_BOX_W = Cm(8.2)
    MID_L       = left_l + LEFT_W + Cm(0.4)
    MID_W       = SW - MID_L - RIGHT_BOX_W - Cm(0.8)
    right_l     = MID_L + MID_W + Cm(0.3)

    # ══════════════════════════════════════════════════════════
    # 左侧栏
    # ══════════════════════════════════════════════════════════
    _rect(sl, left_l, BODY_T2, LEFT_W, PANEL_H, fill=C_LGRAY)

    # 头像
    AV   = Cm(2.8)
    av_l = left_l + (LEFT_W - AV) / 2
    av_t = BODY_T2 + Cm(0.4)
    av_buf = av_cache.get(login)
    placed = False
    if av_buf:
        try:
            av_buf.seek(0)
            sl.shapes.add_picture(av_buf, int(av_l), int(av_t), int(AV), int(AV))
            placed = True
        except Exception:
            pass
    if not placed:
        _rect(sl, av_l, av_t, AV, AV, fill=C_RED)
        _txt(sl, login[:1].upper(),
             av_l, av_t + AV * 0.15, AV, AV * 0.7,
             sz=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 姓名 / @login（固定在头像正下方）
    name_t = av_t + AV + Cm(0.25)
    _txt(sl, name, left_l + Cm(0.3), name_t, LEFT_W - Cm(0.5), Cm(0.85),
         sz=14, bold=True, color=C_DGRAY, align=PP_ALIGN.CENTER)
    _txt(sl, f"@{login}", left_l + Cm(0.3), name_t + Cm(0.9), LEFT_W - Cm(0.5), Cm(0.6),
         sz=11, color=C_RED, align=PP_ALIGN.CENTER)

    # 游标从 @login 下方开始
    LW = _PanelWriter(sl,
                      left=left_l,
                      top=name_t + Cm(1.6),
                      width=LEFT_W,
                      bottom=PANEL_BTM - Cm(0.75),   # 留出底部链接区
                      pad_l=Cm(0.45), pad_r=Cm(0.45))
    LW.hrule(gap_after=Cm(0.15))

    # ── 匹配分数（若有 AI 数据）────────────────────────────────
    if match_score is not None:
        try:
            score_int = max(0, min(100, int(match_score)))
        except (TypeError, ValueError):
            score_int = 0

        LW.txt("华为匹配度", Cm(0.45), sz=9, color=C_MGRAY)
        LW.bar(score_int / 100, h=Cm(0.32), gap_after=Cm(0.06))

        # 分数文字 + 级别徽章（同行）
        if LW.fits(Cm(0.6)):
            _txt(sl, f"{score_int} / 100",
                 left_l + Cm(0.45), LW.y, Cm(3.8), Cm(0.6),
                 sz=11, bold=True, color=C_DGRAY)
            if match_level:
                lc = {"A": C_RED, "B": C_DRED, "C": C_MGRAY}.get(match_level, C_MGRAY)
                LW.badge(f"{match_level} 级", y_offset=Cm(0.06), h=Cm(0.48), color=lc)
            LW.skip(Cm(0.68))

        LW.hrule(gap_after=Cm(0.12))

    # ── 原始信息行（自动溢出压缩）────────────────────────────────
    # 构建要显示的信息行列表
    info_rows = []
    if region:
        info_rows.append(("📍", region))
    if company:
        info_rows.append(("🏢", company))
    if commits:
        info_rows.append(("💻", f"{int(float(commits)):,} commits"))
    if followers:
        info_rows.append(("👥", f"{int(float(followers)):,} followers"))
    if email:
        info_rows.append(("✉", email))
    if blog:
        info_rows.append(("🌐", blog))
    if twitter:
        info_rows.append(("𝕏", f"@{twitter}"))
    if linkedin:
        info_rows.append(("in", linkedin))

    # 根据剩余空间自动选择行高和字号
    if info_rows:
        avail = LW.remaining
        preferred_h = Cm(0.65)
        row_h = preferred_h if len(info_rows) * preferred_h <= avail else max(
            Cm(0.48), avail / len(info_rows))
        sz_info = 10 if row_h >= Cm(0.6) else (9 if row_h >= Cm(0.52) else 8)

        for icon, val in info_rows:
            if not LW.txt(f"{icon}  {val}", row_h, sz=sz_info, color=C_DGRAY, wrap=False):
                break

    # ── 底部链接区（固定位置）────────────────────────────────────
    gh_url = profile_url or f"https://github.com/{login}"
    _txt(sl, gh_url,
         left_l + Cm(0.4), PANEL_BTM - Cm(0.68),
         LEFT_W - Cm(0.6), Cm(0.6),
         sz=8, color=C_RED, url=gh_url, wrap=False)

    # ══════════════════════════════════════════════════════════
    # 中央区
    # ══════════════════════════════════════════════════════════
    MW = _PanelWriter(sl,
                      left=MID_L,
                      top=BODY_T2 + Cm(0.3),
                      width=MID_W,
                      bottom=PANEL_BTM,
                      pad_l=Cm(0.0), pad_r=Cm(0.0))

    # 技术方向
    if tech_direction:
        MW.txt("技术方向", Cm(0.5), sz=10, bold=True, color=C_MGRAY)
        MW.txt(tech_direction, Cm(0.95), sz=14, bold=True, color=C_DGRAY, wrap=True)

    # 技能 badges（不移动 MW.y，手动绘）
    if key_skills and MW.fits(Cm(0.68)):
        bx = MID_L
        for skill in key_skills[:4]:
            sw2 = max(Cm(2.2), Cm(len(skill) * 0.27 + 0.7))
            if bx + sw2 > MID_L + MID_W:
                break
            _rect(sl, bx, MW.y, sw2, Cm(0.58), fill=C_RED)
            _txt(sl, skill, bx + Cm(0.15), MW.y + Cm(0.06),
                 sw2 - Cm(0.2), Cm(0.5),
                 sz=9, bold=True, color=C_WHITE, wrap=False)
            bx += sw2 + Cm(0.2)
        MW.skip(Cm(0.72))

    MW.hrule(gap_after=Cm(0.18))

    # 贡献摘要
    if contribution_summary:
        MW.txt("贡献摘要", Cm(0.5), sz=10, bold=True, color=C_MGRAY)
        MW.txt(contribution_summary, Cm(1.4), sz=11, italic=True,
               color=C_DGRAY, wrap=True)
        MW.skip(Cm(0.1))

    # 来源项目
    if repos:
        MW.txt("来源项目", Cm(0.5), sz=10, bold=True, color=C_MGRAY)
        for repo_name in repos[:8]:
            short = repo_name.split("/")[-1] if "/" in repo_name else repo_name
            if not MW.txt(f"• {short}", Cm(0.52), sz=10, color=C_RED, wrap=False):
                break

    # ══════════════════════════════════════════════════════════
    # 右侧：华为匹配分析
    # ══════════════════════════════════════════════════════════
    right_t = BODY_T2 + Cm(0.3)
    right_h = PANEL_BTM - right_t

    _rect(sl, right_l, right_t, RIGHT_BOX_W, right_h,
          fill=C_LGRAY, line_rgb=C_BORDER, lw_pt=0.6)
    _rect(sl, right_l, right_t, RIGHT_BOX_W, Cm(0.3), fill=C_RED)

    RW = _PanelWriter(sl,
                      left=right_l,
                      top=right_t + Cm(0.45),
                      width=RIGHT_BOX_W,
                      bottom=right_t + right_h - Cm(0.15),
                      pad_l=Cm(0.3), pad_r=Cm(0.3))

    RW.txt("华为匹配分析", Cm(0.6), sz=12, bold=True, color=C_DGRAY)

    if match_score is not None:
        try:
            score_int2 = max(0, min(100, int(match_score)))
        except (TypeError, ValueError):
            score_int2 = 0

        # 大分数 + 级别徽章
        if RW.fits(Cm(1.5)):
            _txt(sl, str(score_int2),
                 right_l + Cm(0.3), RW.y, Cm(2.6), Cm(1.45),
                 sz=40, bold=True, color=C_RED)
            _txt(sl, "/ 100",
                 right_l + Cm(2.2), RW.y + Cm(0.75), Cm(1.6), Cm(0.6),
                 sz=10, color=C_MGRAY)
            if match_level:
                lc2 = {"A": C_RED, "B": C_DRED, "C": C_MGRAY}.get(match_level, C_MGRAY)
                RW.badge(f"{match_level} 级推荐", y_offset=Cm(0.3), h=Cm(0.65), color=lc2)
            RW.skip(Cm(1.58))

        RW.hrule(gap_after=Cm(0.15))

    if match_reason:
        RW.txt("匹配分析", Cm(0.5), sz=10, bold=True, color=C_MGRAY)
        RW.txt(match_reason, RW.remaining, sz=10, italic=True,
               color=C_DGRAY, wrap=True)
    elif match_score is None:
        RW.txt("AI 分析待生成", Cm(0.8), sz=11, color=C_BORDER,
               align=PP_ALIGN.CENTER)
        RW.txt("填写 OpenRouter API Key\n后点击「生成人才档案」",
               Cm(1.2), sz=9, color=C_MGRAY, align=PP_ALIGN.CENTER, wrap=True)


# ════════════════════════════════════════════════════════════
# Slide 6：项目索引页
# ════════════════════════════════════════════════════════════

def slide_project_index(prs, repos_info: dict, talent_map: dict, pg: int, total: int,
                         page_num: int = 1, total_pages_this: int = 1):
    """
    repos_info : {repo: {"tech_areas": [...], ...}}
    talent_map : {repo: [login, ...]}
    """
    sl = prs.slides.add_slide(_blank_layout(prs))
    _rect(sl, 0, 0, SW, SH, fill=C_WHITE)

    header_sub = f"第 {page_num}/{total_pages_this} 页" if total_pages_this > 1 else ""
    _header(sl, "项目贡献者索引", header_sub)
    _footer(sl, f"{pg} / {total}")

    ROW_PAD = Cm(0.2)
    row_t   = BODY_T + Cm(0.4)
    LEFT_W2 = Cm(12)
    right_l2 = Cm(0.5) + LEFT_W2 + Cm(0.25)
    right_w2 = SW - right_l2 - Cm(0.5)
    # Approximate chars that fit per line in right area at sz=9
    CHARS_PER_LINE = 68
    LINE_H = Cm(0.38)

    cy = row_t
    for i, (repo, info) in enumerate(repos_info.items()):
        logins = talent_map.get(repo) or []
        login_str = "  ".join(f"@{l}" for l in logins)

        # Estimate text height for login list
        n_chars = max(1, len(login_str))
        n_lines = max(1, math.ceil(n_chars / CHARS_PER_LINE))
        text_h = n_lines * LINE_H + Cm(0.6)  # 0.3 top + 0.3 bottom pad

        # ROW_H fits both left content (name + badges ≈ 1.35 cm) and right text
        ROW_H = max(Cm(1.45), text_h)

        if cy + ROW_H > BODY_B:
            break  # no more room on this slide

        _rect(sl, Cm(0.5), cy, SW - Cm(1.0), ROW_H,
              fill=C_LGRAY, line_rgb=C_BORDER, lw_pt=0.4)

        # 左区：项目名 + tech_areas
        short = repo.split("/")[-1] if "/" in repo else repo
        _txt(sl, short,
             Cm(0.8), cy + Cm(0.15), LEFT_W2 - Cm(1.0), Cm(0.65),
             sz=11, bold=True, color=C_DGRAY)

        bx = Cm(0.8)
        by = cy + Cm(0.85)
        for tag in (info.get("tech_areas") or [])[:3]:
            tw2 = max(Cm(1.5), Cm(len(tag) * 0.2 + 0.4))
            _rect(sl, bx, by, tw2, Cm(0.38), fill=C_RED)
            _txt(sl, tag,
                 bx + Cm(0.1), by + Cm(0.02), tw2 - Cm(0.15), Cm(0.34),
                 sz=6.5, bold=True, color=C_WHITE, wrap=False)
            bx += tw2 + Cm(0.12)

        # 竖分割线
        _rect(sl, Cm(0.5) + LEFT_W2, cy + Cm(0.15), Cm(0.06), ROW_H - Cm(0.3),
              fill=C_BORDER)

        # 右区：贡献人才 @login 列表（显示全部）
        _txt(sl, login_str or "—",
             right_l2, cy + Cm(0.3), right_w2, ROW_H - Cm(0.5),
             sz=9, color=C_DGRAY, wrap=True)

        cy += ROW_H + ROW_PAD


# ════════════════════════════════════════════════════════════
# 主构建函数
# ════════════════════════════════════════════════════════════

def build_insight_ppt(regions_talents: dict, repos_info: dict,
                       llm_content: dict, theme: str = "华为经典",
                       report_title: str = "开源人才洞察报告",
                       template_bytes: bytes = None) -> io.BytesIO:
    """
    构建洞察报告 PPT。

    参数
    ----
    regions_talents : {region: [talent_dict, ...]}
                      每个 talent_dict 需含 _repos: list[str]
    repos_info      : {repo: {"description", "language", "stars"}}
    llm_content     : {"profiles": {login: {...}}, "overview": {...}}
    theme           : 主题名称
    report_title    : 报告标题

    返回
    ----
    io.BytesIO (PPT 文件内容)
    """
    _apply_theme(theme)

    if template_bytes:
        prs = Presentation(io.BytesIO(template_bytes))
        # 清空模板中已有的幻灯片（保留幻灯片母版）
        xml_slides = prs.slides._sldIdLst
        for sId in list(xml_slides):
            xml_slides.remove(sId)
        prs.slide_width  = SW
        prs.slide_height = SH
    else:
        prs = Presentation()
        prs.slide_width  = SW
        prs.slide_height = SH

    profiles = (llm_content or {}).get("profiles") or {}
    overview = (llm_content or {}).get("overview") or {}

    # 将 AI 生成的 project_tech_map 合并进 repos_info
    project_tech_map = overview.get("project_tech_map") or {}
    enriched_repos: dict = {}
    for repo, info in repos_info.items():
        enriched = dict(info)
        if repo in project_tech_map:
            enriched.update(project_tech_map[repo])
        enriched_repos[repo] = enriched

    # 预计算总页数
    all_talents = [t for ts in regions_talents.values() for t in ts]
    n_regions       = len(regions_talents)
    n_talents       = len(all_talents)
    n_repos         = len(repos_info)
    n_index_pages   = max(1, math.ceil(n_repos / 6))
    n_huawei_pages  = max(1, math.ceil(n_repos / 6))
    total_pages = 1 + 1 + n_huawei_pages + n_regions + n_talents + n_index_pages

    # 预拉取头像（跨地区去重）
    av_cache: dict = {}
    for t in all_talents:
        login = str(t.get("login") or "")
        if login and login not in av_cache:
            av_cache[login] = _fetch_avatar(t.get("avatar_url"))

    pg = 1
    today = datetime.date.today().strftime("%Y  /  %m  /  %d")

    # 预计算匹配度分布
    match_distribution: dict = {}
    for t in all_talents:
        login_t = str(t.get("login") or "")
        level = (profiles.get(login_t) or {}).get("match_level") or ""
        if level:
            match_distribution[level] = match_distribution.get(level, 0) + 1

    # 1. 封面
    slide_cover(prs, report_title, "开源人才洞察报告", today)
    pg += 1

    # 2. 总览 - 质量与密度
    talent_count_by_region = {r: len(ts) for r, ts in regions_talents.items()}
    slide_overview_quality(
        prs,
        quality_summary=overview.get("quality_summary") or "",
        density_stats=overview.get("density_stats") or "",
        talent_count_by_region=talent_count_by_region,
        pg=pg, total=total_pages,
        match_distribution=match_distribution,
        total_count=len(all_talents),
    )
    pg += 1

    # 3. 总览 - 华为匹配（分页显示所有项目，每页6个）
    repo_items_huawei = list(enriched_repos.items())
    for h_idx in range(n_huawei_pages):
        chunk_huawei = dict(repo_items_huawei[h_idx * 6: (h_idx + 1) * 6])
        slide_overview_huawei_placeholder(
            prs, chunk_huawei, pg=pg, total=total_pages,
            page_num=h_idx + 1, total_pages_this=n_huawei_pages,
        )
        pg += 1

    # 4. 各地区：分节页 + 人才详情（按匹配等级/分数排序）
    _LEVEL_ORDER = {"A": 0, "B": 1, "C": 2, "D": 3}
    for region, talents in regions_talents.items():
        sorted_talents = sorted(
            talents,
            key=lambda t: (
                _LEVEL_ORDER.get(
                    (profiles.get(str(t.get("login") or "")) or {}).get("match_level") or "",
                    4,
                ),
                -((profiles.get(str(t.get("login") or "")) or {}).get("match_score") or 0),
            ),
        )
        slide_region_divider(prs, region, len(sorted_talents), pg=pg, total=total_pages)
        pg += 1

        for talent in sorted_talents:
            login = str(talent.get("login") or "")
            talent_data = dict(talent)
            talent_data["_ai"] = profiles.get(login) or {}
            slide_talent(
                prs, talent_data,
                talent.get("_repos") or [],
                pg=pg, total=total_pages,
                av_cache=av_cache,
            )
            pg += 1

    # 5. 项目索引页
    talent_map: dict = {}
    for talent in all_talents:
        login = str(talent.get("login") or "")
        for repo in (talent.get("_repos") or []):
            talent_map.setdefault(repo, [])
            if login not in talent_map[repo]:
                talent_map[repo].append(login)

    repo_items  = list(enriched_repos.items())
    n_index     = max(1, math.ceil(len(repo_items) / 6))
    for idx in range(n_index):
        chunk = dict(repo_items[idx * 6: (idx + 1) * 6])
        slide_project_index(
            prs, chunk, talent_map,
            pg=pg, total=total_pages,
            page_num=idx + 1, total_pages_this=n_index,
        )
        pg += 1

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
